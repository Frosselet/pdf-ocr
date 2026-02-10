"""PowerPoint (PPTX/PPT) table extraction.

Extracts structured tables from PowerPoint presentations.
- PPTX (Office 2007+): Uses python-pptx
- PPT (Office 97-2003): Converts via LibreOffice if available

PPTX has two sources of tabular data:

1. **Table shapes**: Explicit tables with grid structure (easy)
2. **Text shapes**: Text boxes arranged in table-like patterns (requires geometry heuristics)

For text shapes, we apply geometry heuristics:
- H2: Y-clustering to detect rows
- H9: X-position alignment for column unification

Semantic heuristics from heuristics.py apply for header detection and type patterns.
"""

from __future__ import annotations

import shutil
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING

from pdf_ocr.heuristics import (
    StructuredTable,
    TableMetadata,
    build_column_names_from_headers,
    estimate_header_rows,
    normalize_grid,
    split_headers_from_data,
)

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide
    from pptx.table import Table, _Cell


# ---------------------------------------------------------------------------
# Visual element structures
# ---------------------------------------------------------------------------


@dataclass
class PptxCellStyle:
    """Visual style information for a table cell."""

    fill_color: tuple[float, float, float] | None = None
    is_bold: bool = False
    is_italic: bool = False
    font_size: float | None = None


@dataclass
class PptxVisualInfo:
    """Visual information for the entire table."""

    header_fill_rows: list[int] = field(default_factory=list)


@dataclass
class PptxTable:
    """A table extracted from a PowerPoint slide."""

    grid: list[list[str]]
    styles: list[list[PptxCellStyle]] | None = None
    visual: PptxVisualInfo | None = None
    slide_number: int = 0
    table_index: int = 0
    from_shapes: bool = False  # True if reconstructed from text shapes


# ---------------------------------------------------------------------------
# Color utilities
# ---------------------------------------------------------------------------


def _rgb_to_tuple(rgb) -> tuple[float, float, float] | None:
    """Convert python-pptx RGB color to tuple."""
    if rgb is None:
        return None
    try:
        # RGBColor has red, green, blue as 0-255 values
        if hasattr(rgb, "__iter__") and len(rgb) == 3:
            return (rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0)
        # Sometimes it's an RGB object
        r = getattr(rgb, "red", None) or getattr(rgb, "__getitem__", lambda x: 0)(0)
        g = getattr(rgb, "green", None) or getattr(rgb, "__getitem__", lambda x: 0)(1)
        b = getattr(rgb, "blue", None) or getattr(rgb, "__getitem__", lambda x: 0)(2)
        if r is not None and g is not None and b is not None:
            # Could be 0-255 or 0-1
            if max(r, g, b) > 1:
                return (r / 255.0, g / 255.0, b / 255.0)
            return (float(r), float(g), float(b))
    except Exception:
        pass
    return None


def _extract_cell_style(cell: "_Cell") -> PptxCellStyle:
    """Extract visual style from a PPTX table cell."""
    style = PptxCellStyle()

    try:
        # Fill color
        fill = cell.fill
        if fill and fill.type is not None:
            fg_color = getattr(fill.fore_color, "rgb", None)
            if fg_color:
                style.fill_color = _rgb_to_tuple(fg_color)
    except Exception:
        pass

    # Font properties from first paragraph
    try:
        tf = cell.text_frame
        if tf and tf.paragraphs:
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.bold:
                        style.is_bold = True
                    if run.font.italic:
                        style.is_italic = True
                    if run.font.size:
                        style.font_size = run.font.size.pt
                    break  # Just check first run
                break  # Just check first paragraph
    except Exception:
        pass

    return style


# ---------------------------------------------------------------------------
# Table shape extraction
# ---------------------------------------------------------------------------


def _extract_table_shape(
    table: "Table",
    slide_number: int,
    table_index: int,
) -> PptxTable:
    """Extract grid from a PPTX table shape."""
    grid: list[list[str]] = []
    styles: list[list[PptxCellStyle]] = []

    for row in table.rows:
        row_data: list[str] = []
        row_styles: list[PptxCellStyle] = []

        for cell in row.cells:
            # Get cell text
            text = cell.text.strip() if cell.text else ""
            row_data.append(text)

            # Get style
            row_styles.append(_extract_cell_style(cell))

        grid.append(row_data)
        styles.append(row_styles)

    # Analyze visual structure
    visual = _analyze_visual_structure(styles)

    return PptxTable(
        grid=grid,
        styles=styles,
        visual=visual,
        slide_number=slide_number,
        table_index=table_index,
        from_shapes=False,
    )


# ---------------------------------------------------------------------------
# Text shape clustering (geometry heuristics)
# ---------------------------------------------------------------------------


@dataclass
class TextSpan:
    """A text span with position information."""

    text: str
    x: float  # Left position
    y: float  # Top position
    width: float
    height: float
    is_bold: bool = False


def _extract_text_shapes(slide: "Slide") -> list[TextSpan]:
    """Extract text spans from all text shapes on a slide."""
    spans: list[TextSpan] = []

    for shape in slide.shapes:
        # Skip table shapes (handled separately)
        if shape.has_table:
            continue

        # Get text from text frames
        if not shape.has_text_frame:
            continue

        try:
            tf = shape.text_frame
            text = tf.text.strip() if tf.text else ""
            if not text:
                continue

            # Get position (EMUs to points: 1 point = 914400 EMUs)
            emu_to_pt = 914400
            x = shape.left / emu_to_pt if shape.left else 0
            y = shape.top / emu_to_pt if shape.top else 0
            width = shape.width / emu_to_pt if shape.width else 0
            height = shape.height / emu_to_pt if shape.height else 0

            # Check for bold
            is_bold = False
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.bold:
                        is_bold = True
                        break
                if is_bold:
                    break

            spans.append(TextSpan(
                text=text,
                x=x,
                y=y,
                width=width,
                height=height,
                is_bold=is_bold,
            ))
        except Exception:
            continue

    return spans


def _cluster_y_positions(spans: list[TextSpan], threshold: float = 5.0) -> list[list[TextSpan]]:
    """H2: Cluster spans by y-position to form rows.

    Spans with y-positions within threshold of each other are grouped.
    """
    if not spans:
        return []

    # Sort by y position
    sorted_spans = sorted(spans, key=lambda s: s.y)

    # Greedy clustering
    rows: list[list[TextSpan]] = [[sorted_spans[0]]]

    for span in sorted_spans[1:]:
        # Check if span belongs to current row
        current_row_y = sum(s.y for s in rows[-1]) / len(rows[-1])
        if abs(span.y - current_row_y) <= threshold:
            rows[-1].append(span)
        else:
            rows.append([span])

    return rows


def _unify_columns(rows: list[list[TextSpan]], tolerance: float = 10.0) -> tuple[list[float], dict[float, int]]:
    """H9: Build canonical column positions from all spans.

    Returns (canonical_positions, x_to_col_index).
    """
    # Collect all x positions
    all_x: list[float] = []
    for row in rows:
        for span in row:
            all_x.append(span.x)

    if not all_x:
        return [], {}

    all_x.sort()

    # Cluster x positions
    clusters: list[list[float]] = [[all_x[0]]]
    for x in all_x[1:]:
        cluster_mean = sum(clusters[-1]) / len(clusters[-1])
        if abs(x - cluster_mean) <= tolerance:
            clusters[-1].append(x)
        else:
            clusters.append([x])

    # Build mapping
    canonical: list[float] = []
    x_to_col: dict[float, int] = {}

    for ci, cluster in enumerate(clusters):
        canonical.append(sum(cluster) / len(cluster))
        for x in cluster:
            x_to_col[x] = ci

    return canonical, x_to_col


def _spans_to_grid(
    rows: list[list[TextSpan]],
    canonical: list[float],
    x_to_col: dict[float, int],
) -> list[list[str]]:
    """Convert clustered spans to a grid."""
    num_cols = len(canonical)
    if num_cols == 0:
        return []

    grid: list[list[str]] = []

    for row in rows:
        cells = [""] * num_cols
        row_sorted = sorted(row, key=lambda s: s.x)

        for span in row_sorted:
            # Find column for this span
            # Use exact x if in map, otherwise find closest
            if span.x in x_to_col:
                col_idx = x_to_col[span.x]
            else:
                # Find closest canonical position
                min_dist = float("inf")
                col_idx = 0
                for ci, pos in enumerate(canonical):
                    dist = abs(span.x - pos)
                    if dist < min_dist:
                        min_dist = dist
                        col_idx = ci

            if col_idx < num_cols:
                if cells[col_idx]:
                    cells[col_idx] += " " + span.text
                else:
                    cells[col_idx] = span.text

        grid.append(cells)

    return grid


def _detect_table_from_shapes(
    slide: "Slide",
    slide_number: int,
    min_rows: int = 3,
    min_cols: int = 2,
) -> PptxTable | None:
    """Try to detect a table from arranged text shapes.

    Uses geometry heuristics to find table-like patterns.
    """
    spans = _extract_text_shapes(slide)
    if len(spans) < min_rows * min_cols:
        return None

    # Cluster into rows
    rows = _cluster_y_positions(spans)
    if len(rows) < min_rows:
        return None

    # Unify columns
    canonical, x_to_col = _unify_columns(rows)
    if len(canonical) < min_cols:
        return None

    # Check consistency - most rows should have similar span counts
    span_counts = [len(row) for row in rows]
    if max(span_counts) - min(span_counts) > 3:
        # Too variable, probably not a table
        return None

    # Convert to grid
    grid = _spans_to_grid(rows, canonical, x_to_col)

    return PptxTable(
        grid=grid,
        styles=None,  # No style info from text shapes
        visual=None,
        slide_number=slide_number,
        table_index=0,
        from_shapes=True,
    )


# ---------------------------------------------------------------------------
# Visual analysis
# ---------------------------------------------------------------------------


def _analyze_visual_structure(
    styles: list[list[PptxCellStyle]],
) -> PptxVisualInfo:
    """Analyze visual structure of the table."""
    visual = PptxVisualInfo()

    if not styles or not styles[0]:
        return visual

    # Check for header fills in first few rows
    row_fills: list[tuple[float, float, float] | None] = []
    for ri, row in enumerate(styles[:5]):
        fills = [s.fill_color for s in row if s.fill_color is not None]
        if fills:
            row_fills.append(fills[0])
        else:
            row_fills.append(None)

    # Find header rows with consistent fill
    if row_fills and row_fills[0] is not None:
        first_fill = row_fills[0]
        for ri, fill in enumerate(row_fills):
            if fill == first_fill:
                visual.header_fill_rows.append(ri)
            else:
                break

    return visual


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def extract_tables_from_pptx(
    pptx_path: str | Path,
    *,
    slides: list[int] | None = None,
    include_shape_tables: bool = True,
    extract_styles: bool = True,
) -> list[StructuredTable]:
    """Extract structured tables from a PowerPoint presentation.

    Args:
        pptx_path: Path to the .pptx file
        slides: Optional list of 0-based slide indices to process.
            If None, processes all slides.
        include_shape_tables: If True, also detect tables from text shapes.
        extract_styles: If True, extract visual style information.

    Returns:
        List of StructuredTable objects, one per table found.

    Raises:
        ImportError: If python-pptx is not installed.
        FileNotFoundError: If the file doesn't exist.
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError(
            "python-pptx is required for PPTX extraction. "
            "Install with: pip install pdf-ocr[pptx]"
        ) from e

    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs: Presentation = Presentation(path)
    result: list[StructuredTable] = []

    # Determine which slides to process
    slide_indices = slides if slides is not None else range(len(prs.slides))

    for slide_idx in slide_indices:
        if slide_idx >= len(prs.slides):
            continue

        slide = prs.slides[slide_idx]
        table_count = 0

        # Extract table shapes
        for shape in slide.shapes:
            if shape.has_table:
                pptx_table = _extract_table_shape(
                    shape.table,
                    slide_number=slide_idx,
                    table_index=table_count,
                )
                table_count += 1

                # Process table
                grid = normalize_grid(pptx_table.grid)
                if not grid or not grid[0]:
                    continue

                # Estimate headers
                header_count = 0
                if pptx_table.visual and pptx_table.visual.header_fill_rows:
                    header_count = len(pptx_table.visual.header_fill_rows)
                else:
                    header_count = estimate_header_rows(grid)

                header_rows, data_rows = split_headers_from_data(grid, header_count)
                column_names = build_column_names_from_headers(header_rows)

                # Ensure correct column count
                num_cols = len(grid[0]) if grid else 0
                while len(column_names) < num_cols:
                    column_names.append("")

                result.append(StructuredTable(
                    column_names=column_names,
                    data=data_rows,
                    source_format="pptx",
                    metadata={
                        "slide_number": slide_idx,
                        "table_index": pptx_table.table_index,
                        "from_shapes": False,
                    },
                ))

        # Try to detect table from text shapes
        if include_shape_tables:
            shape_table = _detect_table_from_shapes(slide, slide_idx)
            if shape_table is not None:
                grid = normalize_grid(shape_table.grid)
                if grid and grid[0]:
                    header_count = estimate_header_rows(grid)
                    header_rows, data_rows = split_headers_from_data(grid, header_count)
                    column_names = build_column_names_from_headers(header_rows)

                    num_cols = len(grid[0]) if grid else 0
                    while len(column_names) < num_cols:
                        column_names.append("")

                    result.append(StructuredTable(
                        column_names=column_names,
                        data=data_rows,
                        source_format="pptx",
                        metadata={
                            "slide_number": slide_idx,
                            "table_index": table_count,
                            "from_shapes": True,
                        },
                    ))

    return result


# ---------------------------------------------------------------------------
# PPT (Legacy PowerPoint 97-2003) extraction via conversion
# ---------------------------------------------------------------------------


def _find_libreoffice() -> str | None:
    """Find LibreOffice executable on the system."""
    candidates = [
        "libreoffice",
        "soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
        "C:\\Program Files\\LibreOffice\\program\\soffice.exe",  # Windows
    ]

    for candidate in candidates:
        if shutil.which(candidate):
            return candidate

    return None


def _convert_ppt_to_pptx(ppt_path: Path, output_dir: Path) -> Path | None:
    """Convert a .ppt file to .pptx using LibreOffice."""
    libreoffice = _find_libreoffice()
    if not libreoffice:
        return None

    try:
        result = subprocess.run(
            [
                libreoffice,
                "--headless",
                "--convert-to", "pptx",
                "--outdir", str(output_dir),
                str(ppt_path),
            ],
            capture_output=True,
            timeout=60,
        )

        if result.returncode != 0:
            return None

        expected_name = ppt_path.stem + ".pptx"
        converted_path = output_dir / expected_name

        if converted_path.exists():
            return converted_path

        return None

    except (subprocess.TimeoutExpired, subprocess.SubprocessError):
        return None


def _extract_tables_from_ppt(
    ppt_path: str | Path,
    *,
    slides: list[int] | None = None,
    include_shape_tables: bool = True,
    extract_styles: bool = True,
) -> list[StructuredTable]:
    """Extract tables from a legacy .ppt file by converting to .pptx first."""
    path = Path(ppt_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {ppt_path}")

    if not _find_libreoffice():
        raise ImportError(
            "LibreOffice is required to read legacy .ppt files. "
            "Please install LibreOffice or convert the file to .pptx format. "
            "Download: https://www.libreoffice.org/download/"
        )

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)
        converted = _convert_ppt_to_pptx(path, tmpdir_path)

        if not converted:
            raise RuntimeError(
                f"Failed to convert {path.name} to .pptx format. "
                "Please convert the file manually or check your LibreOffice installation."
            )

        result = extract_tables_from_pptx(
            converted,
            slides=slides,
            include_shape_tables=include_shape_tables,
            extract_styles=extract_styles,
        )

        for table in result:
            table.source_format = "ppt"

        return result


def extract_tables_from_powerpoint(
    pptx_path: str | Path,
    *,
    slides: list[int] | None = None,
    include_shape_tables: bool = True,
    extract_styles: bool = True,
) -> list[StructuredTable]:
    """Extract structured tables from a PowerPoint file (PPT or PPTX).

    Auto-detects format based on file extension.

    Args:
        pptx_path: Path to the PowerPoint file (.pptx or .ppt)
        slides: Optional list of 0-based slide indices to process.
        include_shape_tables: If True, also detect tables from text shapes.
        extract_styles: If True, extract visual style information.

    Returns:
        List of StructuredTable objects.

    Note: .ppt files require LibreOffice for conversion.
    """
    path = Path(pptx_path)
    suffix = path.suffix.lower()

    if suffix == ".pptx":
        return extract_tables_from_pptx(
            path, slides=slides, include_shape_tables=include_shape_tables,
            extract_styles=extract_styles,
        )
    elif suffix == ".ppt":
        return _extract_tables_from_ppt(
            path, slides=slides, include_shape_tables=include_shape_tables,
            extract_styles=extract_styles,
        )
    else:
        try:
            return extract_tables_from_pptx(
                path, slides=slides, include_shape_tables=include_shape_tables,
                extract_styles=extract_styles,
            )
        except Exception:
            return _extract_tables_from_ppt(
                path, slides=slides, include_shape_tables=include_shape_tables,
                extract_styles=extract_styles,
            )


def pptx_to_markdown(
    pptx_path: str | Path,
    *,
    slides: list[int] | None = None,
) -> str:
    """Extract tables from PowerPoint and render as markdown.

    Convenience function for quick inspection of presentation content.
    """
    tables = extract_tables_from_pptx(pptx_path, slides=slides, extract_styles=False)

    parts: list[str] = []
    for table in tables:
        meta = table.metadata or {}
        slide = meta.get("slide_number", 0)
        from_shapes = meta.get("from_shapes", False)
        source = " (from shapes)" if from_shapes else ""
        parts.append(f"## Slide {slide + 1}{source}\n")

        if table.column_names:
            parts.append("| " + " | ".join(table.column_names) + " |")
            parts.append("| " + " | ".join(["---"] * len(table.column_names)) + " |")

        for row in table.data:
            parts.append("| " + " | ".join(row) + " |")

        parts.append("")  # Blank line between tables

    return "\n".join(parts)
